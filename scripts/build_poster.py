#!/usr/bin/env python3
"""
Generate the team046 project poster as a single-slide 30x40 inch .pptx.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(REPO, "docs", "team046poster.pptx")

W_IN, H_IN = 30, 40
W, H = Inches(W_IN), Inches(H_IN)

# palette
CRIMSON     = RGBColor(0xB7, 0x1C, 0x1C)
DARK_RED    = RGBColor(0x8B, 0x0A, 0x0A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xFA, 0xFA, 0xFA)
NEAR_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
BODY_TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
CAPTION_CLR = RGBColor(0x66, 0x66, 0x66)
SECTION_BG  = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_LINE = RGBColor(0xCC, 0x33, 0x33)

MARGIN = 0.5
GAP = 0.35
COL_GAP = 0.4
COL_W = (W_IN - 2 * MARGIN - COL_GAP) / 2  # ~14.55 each
LEFT_X = MARGIN
RIGHT_X = MARGIN + COL_W + COL_GAP
INNER_PAD = 0.35


def _box(slide, x, y, w, h, fill=SECTION_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.02
    return shape


def _rect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _tb(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def _heading(slide, x, y, w, text):
    # accent bar
    _rect(slide, x, y, 0.12, 0.5, CRIMSON)
    tb = _tb(slide, x + 0.22, y, w - 0.22, 0.55)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CRIMSON
    p.font.name = "Calibri"
    return y + 0.65


def _bullets(slide, x, y, w, h, items, size=22):
    tb = _tb(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
        p.text = f"\u2022   {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = BODY_TEXT
        p.font.name = "Calibri"
        p.line_spacing = Pt(size + 8)


def _caption(slide, x, y, w, text):
    tb = _tb(slide, x, y, w, 0.4)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = CAPTION_CLR
    p.font.name = "Calibri"
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER


def _img(slide, rel_path, x, y, width=None, height=None):
    full = os.path.join(REPO, rel_path)
    if not os.path.exists(full):
        print(f"  [warn] missing: {full}")
        return None
    kw = {"left": Inches(x), "top": Inches(y)}
    if width:
        kw["width"] = Inches(width)
    if height:
        kw["height"] = Inches(height)
    return slide.shapes.add_picture(full, **kw)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # ═══════════════════════════════════════════════════════════
    # TITLE BANNER
    # ═══════════════════════════════════════════════════════════
    _rect(slide, 0, 0, W_IN, 3.4, CRIMSON)
    _rect(slide, 0, 3.2, W_IN, 0.2, DARK_RED)

    tb = _tb(slide, 1, 0.45, W_IN - 2, 1.3)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Ignition Insights"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    tb = _tb(slide, 1, 1.6, W_IN - 2, 0.7)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Explaining and Forecasting Urban Fire Risk"
    p.font.size = Pt(36)
    p.font.bold = False
    p.font.color.rgb = RGBColor(0xFF, 0xDD, 0xDD)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    tb = _tb(slide, 1, 2.35, W_IN - 2, 0.7)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Vishruth Anand   \u00b7   Vineeth Nareddy   \u00b7   Rian Rahman   \u00b7   James Reilly   \u00b7   Jayanth Vennamreddy"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Georgia Institute of Technology  \u2014  CSE 6242 Data and Visual Analytics  \u2014  Team 46"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xFF, 0xBB, 0xBB)
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    Y0 = 3.8
    pad = INNER_PAD
    lx = LEFT_X
    rx = RIGHT_X
    cw = COL_W

    # ═══════════════════════════════════════════════════════════
    # LEFT COL: INTRODUCTION
    # ═══════════════════════════════════════════════════════════
    sec_h = 5.8
    _box(slide, lx, Y0, cw, sec_h)
    by = _heading(slide, lx + pad, Y0 + pad, cw - 2*pad, "Introduction")
    _bullets(slide, lx + pad, by, cw - 2*pad, sec_h - 1.4, [
        "Fire departments make staffing and inspection decisions\n"
        "under uncertainty, yet fire incidents cluster unevenly\n"
        "across neighborhoods and time periods.",
        "Current tools show what happened \u2014 not where risk is\n"
        "concentrated, how it shifts over time, or why a model\n"
        "flags a particular area.",
        "Better risk awareness helps departments pre-position\n"
        "resources and prioritize inspections, potentially\n"
        "reducing response times and preventing incidents.",
        "We built an end-to-end system that forecasts next-day\n"
        "fire risk on a 1 km Atlanta grid and explains each\n"
        "prediction through an interactive dashboard.",
    ], size=22)

    # ═══════════════════════════════════════════════════════════
    # LEFT COL: DATA
    # ═══════════════════════════════════════════════════════════
    y2 = Y0 + sec_h + GAP
    sec_h2 = 4.2
    _box(slide, lx, y2, cw, sec_h2)
    by = _heading(slide, lx + pad, y2 + pad, cw - 2*pad, "Data")
    _bullets(slide, lx + pad, by, cw - 2*pad, sec_h2 - 1.4, [
        "1,473 geocoded fire incidents from 2024 NFIRS PDR Light\n"
        "(Atlanta, fire codes 100\u2013199, deduplicated).",
        "Daily weather from Open-Meteo: temperature, humidity,\n"
        "precipitation, wind speed (citywide, one row per day).",
        "1,176 grid cells at 1 km resolution; 132,860 cell-day rows\n"
        "with 18 features (weather, calendar, lags, rolling sums).",
    ], size=22)

    # ═══════════════════════════════════════════════════════════
    # LEFT COL: METHODS
    # ═══════════════════════════════════════════════════════════
    y3 = y2 + sec_h2 + GAP
    sec_h3 = 8.2
    _box(slide, lx, y3, cw, sec_h3)
    by = _heading(slide, lx + pad, y3 + pad, cw - 2*pad, "Methods")
    _bullets(slide, lx + pad, by, cw - 2*pad, sec_h3 - 1.4, [
        "Hotspot baseline: Gaussian KDE over training centroids\n"
        "weighted by count. Captures static spatial density only.",
        "ARIMA baseline: per-cell time-series for top 50 cells.\n"
        "Grid search (p,d,q) by AIC; fallback for sparse cells.",
        "Random Forest (main model): 300 trees, max depth 15,\n"
        "balanced weights. Uses all 14 features + encoded grid_id.\n"
        "Combines weather, calendar, history, and location signals\n"
        "that simpler baselines ignore individually.",
        "Temporal 80/20 split: train Jan\u2013Oct 18, test Oct 19\u2013Dec 30.\n"
        "No row shuffling to prevent future-data leakage.",
        "Key novelty: most prior systems provide either forecasting\n"
        "or explanation. Ours combines short-horizon prediction,\n"
        "SHAP explanation, and interactive exploration in one tool.",
    ], size=22)

    # ═══════════════════════════════════════════════════════════
    # LEFT COL: RISK HEATMAP
    # ═══════════════════════════════════════════════════════════
    y4 = y3 + sec_h3 + GAP
    fig_h = 8.5
    _box(slide, lx, y4, cw, fig_h + 1.0)
    _img(slide, "outputs/maps/risk_map_aggregate.png",
         lx + (cw - 8.0) / 2, y4 + 0.3, width=8.0)
    _caption(slide, lx + pad, y4 + fig_h + 0.15, cw - 2*pad,
             "Fig. 1: Aggregate RF risk across all test dates (Atlanta 1 km grid)")

    # ═══════════════════════════════════════════════════════════
    # LEFT COL: DASHBOARD
    # ═══════════════════════════════════════════════════════════
    y5 = y4 + fig_h + 1.0 + GAP
    sec_h5 = H_IN - MARGIN - y5
    _box(slide, lx, y5, cw, sec_h5)
    by = _heading(slide, lx + pad, y5 + pad, cw - 2*pad, "Interactive Dashboard")
    _bullets(slide, lx + pad, by, cw - 2*pad, 3.0, [
        "Leaflet map + D3 analytics. Metric selector, date slider\n"
        "with play/pause, linked histogram and time-series views.",
        "Click any cell: feature snapshot, SHAP explanation panel\n"
        "with driver bar chart and narrative sentence.",
    ], size=22)
    _img(slide, "outputs/frontend-captures/app_full.png",
         lx + (cw - 12.5) / 2, by + 3.2, width=12.5)
    _caption(slide, lx + pad, by + 3.2 + 7.8, cw - 2*pad,
             "Fig. 2: Dashboard with choropleth, time control, and explanation panel")

    # ═══════════════════════════════════════════════════════════
    # RIGHT COL: RESULTS
    # ═══════════════════════════════════════════════════════════
    _box(slide, rx, Y0, cw, 4.6)
    by = _heading(slide, rx + pad, Y0 + pad, cw - 2*pad, "Results")
    _bullets(slide, rx + pad, by, cw - 2*pad, 3.6, [
        "Test set: 26,572 cell-days, only 138 positive (0.52%).\n"
        "Extreme imbalance makes accuracy misleading (all > 98%).",
        "RF achieves the best ranking: ROC-AUC 0.65, PR-AUC 0.014\n"
        "(\u22483\u00d7 better than random at this positive rate).",
        "RF threshold gives 0 true positives, but the continuous\n"
        "probability surface is more useful as a risk map.",
    ], size=22)

    # ═══════════════════════════════════════════════════════════
    # RIGHT COL: METRICS FIGURE
    # ═══════════════════════════════════════════════════════════
    ry2 = Y0 + 4.6 + GAP
    _box(slide, rx, ry2, cw, 7.0)
    _img(slide, "baselines/outputs/plots/metrics_comparison.png",
         rx + (cw - 12.0) / 2, ry2 + 0.3, width=12.0)
    _caption(slide, rx + pad, ry2 + 6.2, cw - 2*pad,
             "Fig. 3: Model comparison across key evaluation metrics")

    # ═══════════════════════════════════════════════════════════
    # RIGHT COL: ROC FIGURE
    # ═══════════════════════════════════════════════════════════
    ry3 = ry2 + 7.0 + GAP
    _box(slide, rx, ry3, cw, 8.6)
    _img(slide, "baselines/outputs/plots/roc_curves.png",
         rx + (cw - 9.5) / 2, ry3 + 0.3, width=9.5)
    _caption(slide, rx + pad, ry3 + 7.9, cw - 2*pad,
             "Fig. 4: ROC curves for all three baseline models")

    # ═══════════════════════════════════════════════════════════
    # RIGHT COL: INTERPRETABILITY + SHAP
    # ═══════════════════════════════════════════════════════════
    ry4 = ry3 + 8.6 + GAP
    sec_h_interp = 10.0
    _box(slide, rx, ry4, cw, sec_h_interp)
    by = _heading(slide, rx + pad, ry4 + pad, cw - 2*pad, "Interpretability")
    _bullets(slide, rx + pad, by, cw - 2*pad, 2.0, [
        "SHAP TreeExplainer on the RF model: per-feature\n"
        "contributions for every cell-date prediction.",
        "Top drivers: rolling_sum_14, rolling_sum_7 (recent activity),\n"
        "then month, wind, and grid location.",
    ], size=22)
    _img(slide, "outputs/interpretability/shap_bar.png",
         rx + (cw - 9.0) / 2, by + 2.3, width=9.0)
    _caption(slide, rx + pad, by + 2.3 + 7.0, cw - 2*pad,
             "Fig. 5: SHAP feature importance (mean |SHAP value|)")

    # ═══════════════════════════════════════════════════════════
    # RIGHT COL: CONCLUSIONS
    # ═══════════════════════════════════════════════════════════
    ry5 = ry4 + sec_h_interp + GAP
    sec_h_conc = 3.8
    _box(slide, rx, ry5, cw, sec_h_conc)
    by = _heading(slide, rx + pad, ry5 + pad, cw - 2*pad, "Conclusions")
    _bullets(slide, rx + pad, by, cw - 2*pad, sec_h_conc - 1.2, [
        "Binary metrics are insufficient at 0.5% positive rate \u2014\n"
        "ranked risk surfaces are more practical than alarms.",
        "Linked views (map + timeline + histogram + explanation)\n"
        "make model outputs easier to understand and critique.",
        "Future: census/land-use features, calibrated thresholds,\n"
        "gradient boosting, and a formal user study.",
    ], size=22)

    # ═══════════════════════════════════════════════════════════
    # RIGHT COL: REFERENCES
    # ═══════════════════════════════════════════════════════════
    ry6 = ry5 + sec_h_conc + GAP
    remaining = H_IN - MARGIN - ry6
    _box(slide, rx, ry6, cw, remaining)
    by = _heading(slide, rx + pad, ry6 + pad, cw - 2*pad, "References")
    refs = [
        "[1] Wang et al., CityGuard, IMWUT 2019",
        "[2] Coffield et al., Fire size ML, IJWF 2019",
        "[3] Lattimer et al., ML fire simulation, FSJ 2020",
        "[4] Asgary et al., Toronto fire clustering, 2010",
        "[5] Yuan & Wylie, ARIMA vs RF, Austin 2024",
        "[6] Madaio et al., Firebird Atlanta, 2016",
        "[7] Jin et al., Deep sequence fire, ASC 2020",
        "[8] Ahn et al., Stacking ensemble, Fire 2024",
        "[9\u201315] Zhang, Ku, Cui, Liao, Kang, Xiao, Jennings",
    ]
    _bullets(slide, rx + pad, by, cw - 2*pad, remaining - 1.0, refs, size=16)

    # ═══════════════════════════════════════════════════════════
    # FOOTER BAR
    # ═══════════════════════════════════════════════════════════
    _rect(slide, 0, H_IN - 0.18, W_IN, 0.18, CRIMSON)

    # save
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Saved poster to {OUT}")
    print(f"  Size: {W_IN} x {H_IN} inches")


if __name__ == "__main__":
    build()
